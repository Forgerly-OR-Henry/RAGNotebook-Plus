"""
模块职责：测试模块，使用单元测试和回归用例验证当前业务契约。

主要协作：本文件只声明当前模块的职责边界，运行时行为由下方函数、类和依赖对象共同完成。
"""

import io

import pytest

from agent.indexing.document_parser import DocumentParser
from mvc.services.document_preview_service import render_docx_html, render_pptx_html
from utils.file_handler import ppt_loader_sync, word_loader_sync


def test_word_loader_sync_extracts_docx_without_unstructured(tmp_path):
    """
    用途：执行test word loader sync extracts docx without unstructured相关业务逻辑。

    参数：
    - tmp_path（未显式标注）：调用方传入的tmp_path数据或控制参数，用于驱动本函数处理流程。

    返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
    """
    docx = pytest.importorskip("docx")
    document = docx.Document()
    document.add_heading("计算机体系结构", level=1)
    document.add_paragraph("流水线、缓存和指令集是核心内容")

    file_path = tmp_path / "architecture.docx"
    document.save(file_path)

    docs = word_loader_sync(str(file_path))

    assert len(docs) == 1
    assert "计算机体系结构" in docs[0].page_content
    assert "流水线、缓存和指令集" in docs[0].page_content
    assert docs[0].metadata["file_type"] == "docx"


def test_ppt_loader_sync_extracts_pptx_slides(tmp_path):
    """
    用途：执行test ppt loader sync extracts pptx slides相关业务逻辑。

    参数：
    - tmp_path（未显式标注）：调用方传入的tmp_path数据或控制参数，用于驱动本函数处理流程。

    返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
    """
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "课程结构"
    slide.placeholders[1].text = "第一章 CPU\n第二章 存储系统"

    file_path = tmp_path / "course.pptx"
    presentation.save(file_path)

    docs = ppt_loader_sync(str(file_path))

    assert len(docs) == 1
    assert "幻灯片 1" in docs[0].page_content
    assert "课程结构" in docs[0].page_content
    assert "第一章 CPU" in docs[0].page_content
    assert docs[0].metadata["slide_number"] == 1


def test_document_parser_rejects_empty_docx():
    """
    用途：执行test document parser rejects empty docx相关业务逻辑。

    参数：无显式业务参数。

    返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
    """
    docx = pytest.importorskip("docx")
    document = docx.Document()
    buffer = io.BytesIO()
    document.save(buffer)

    parser = DocumentParser()

    with pytest.raises(ValueError, match="未解析到可索引文本"):
        parser.parse_bytes_sync(content=buffer.getvalue(), filename="empty.docx", user_id="user-1")


def test_docx_html_preview_keeps_formatting_and_tables():
    """
    用途：执行test docx html preview keeps formatting and tables相关业务逻辑。

    参数：无显式业务参数。

    返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
    """
    docx = pytest.importorskip("docx")
    from docx.shared import Pt, RGBColor

    document = docx.Document()
    paragraph = document.add_paragraph()
    run = paragraph.add_run("红色大字")
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    run.font.size = Pt(22)
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "课程"
    table.cell(0, 1).text = "计算机体系结构"

    buffer = io.BytesIO()
    document.save(buffer)

    html = render_docx_html(buffer.getvalue(), "format.docx").decode("utf-8")

    assert "红色大字" in html
    assert "font-size: 22.00pt" in html
    assert "color: #FF0000" in html
    assert "<table" in html
    assert "计算机体系结构" in html


def test_pptx_html_preview_keeps_text_formatting_and_tables():
    """
    用途：执行test pptx html preview keeps text formatting and tables相关业务逻辑。

    参数：无显式业务参数。

    返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
    """
    pptx = pytest.importorskip("pptx")
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "蓝色标题"
    run.font.color.rgb = RGBColor(0x00, 0x50, 0xC8)
    run.font.size = Pt(24)
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(5), Inches(1))
    table_shape.table.cell(0, 0).text = "章节"
    table_shape.table.cell(0, 1).text = "缓存"

    buffer = io.BytesIO()
    presentation.save(buffer)

    html = render_pptx_html(buffer.getvalue(), "slides.pptx").decode("utf-8")

    assert "蓝色标题" in html
    assert "font-size: 24.00pt" in html
    assert "color: #0050C8" in html
    assert "<table" in html
    assert "缓存" in html
