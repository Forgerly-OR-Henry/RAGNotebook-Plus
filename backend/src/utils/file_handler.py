import asyncio
import logging
import os

from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
)
from langchain_core.documents import Document

from core.logger_handler import logger
from utils.path_tool import get_abstract_path


class FontBBoxLogFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:
        return "FontBBox from font descriptor" not in record.getMessage()


_font_bbox_filter = FontBBoxLogFilter()
for _logger_name in ("pypdf", "pypdf._reader", "pypdf.generic"):
    logging.getLogger(_logger_name).addFilter(_font_bbox_filter)

async def listdir_allowed_type(path: str, allowed_types: tuple[str]) -> tuple:
    """
    获取指定目录下所有允许的文件类型
    :param path: 目录路径
    :param allowed_types: 允许的文件类型元组
    :return: 符合条件的文件路径列表
    """
    # 处理路径，确保使用绝对路径
    abs_path = get_abstract_path(path) if not os.path.isabs(path) else path

    if not os.path.exists(abs_path):
        logger.error(f"【文件列表】目录路径 {abs_path} 不存在")
        return ()

    if not os.path.isdir(abs_path):
        logger.error(f"【文件列表】目录路径 {abs_path} 不是目录")
        return ()

    file_list = []
    for f in await asyncio.to_thread(os.listdir, abs_path):
        if f.lower().endswith(allowed_types):
            file_path = os.path.join(abs_path, f)
            file_list.append(file_path)

    return tuple(file_list)



async def pdf_loader(file_path: str, password: str = None) -> list[Document]:
    """
    加载PDF文件内容（支持包含图片和文字的混合PDF）
    :param file_path: PDF文件路径
    :param password: PDF密码（如果有）
    :return: PDF文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    loader = PyPDFLoader(abs_file_path, password=password) if password else PyPDFLoader(abs_file_path)
    return await asyncio.to_thread(loader.load)


async def txt_loader(file_path: str) -> list[Document]:
    """
    加载TXT文件内容
    :param file_path: TXT文件路径
    :return: TXT文件内容
    """
    # 处理路径，确保使用绝对路径
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    # 使用不同的编码加载文件
    encodings = ['utf-8', 'gbk']
    for encoding in encodings:
        try:
            loader = TextLoader(abs_file_path, encoding=encoding)
            return await asyncio.to_thread(loader.load)
        except Exception as e:
            logger.error(f"【文本文件加载】使用编码 {encoding} 加载文件 {abs_file_path} 时出错: {e}")
            continue
    # 所有编码都失败，返回空列表
    return []

async def word_loader(file_path: str) -> list[Document]:
    """
    加载WORD文件内容
    :param file_path: WORD文件路径
    :return: WORD文件内容
    """
    return await asyncio.to_thread(word_loader_sync, file_path)

async def markdown_loader(file_path: str) -> list[Document]:
    """
    加载Markdown文件内容
    :param file_path: Markdown文件路径
    :return: Markdown文件内容
    """
    return await asyncio.to_thread(markdown_loader_sync, file_path)


async def ppt_loader(file_path: str) -> list[Document]:
    """
    加载PPT/PPTX文件内容
    :param file_path: PPT文件路径
    :return: PPT文件内容
    """
    return await asyncio.to_thread(ppt_loader_sync, file_path)


def pdf_loader_sync(file_path: str, password: str = None) -> list[Document]:
    """
    同步加载PDF文件内容（用于多线程场景，支持包含图片和文字的混合PDF）
    :param file_path: PDF文件路径
    :param password: PDF密码（如果有）
    :return: PDF文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    loader = PyPDFLoader(abs_file_path, password=password) if password else PyPDFLoader(abs_file_path)
    return loader.load()


def txt_loader_sync(file_path: str) -> list[Document]:
    """
    同步加载TXT文件内容（用于多线程场景）
    :param file_path: TXT文件路径
    :return: TXT文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    encodings = ['utf-8', 'gbk']
    for encoding in encodings:
        try:
            loader = TextLoader(abs_file_path, encoding=encoding)
            return loader.load()
        except Exception as e:
            logger.error(f"【文本文件加载】使用编码 {encoding} 加载文件 {abs_file_path} 时出错: {e}")
            continue
    return []


def word_loader_sync(file_path: str) -> list[Document]:
    """
    同步加载WORD文件内容（用于多线程场景）
    :param file_path: WORD文件路径
    :return: WORD文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    lower_path = abs_file_path.lower()
    if lower_path.endswith(".doc") and not lower_path.endswith(".docx"):
        logger.error("【WORD文件加载】旧版 .doc 文件需要先另存为 .docx 后上传")
        return []

    try:
        content = _extract_docx_text(abs_file_path)
        return _single_document(content, abs_file_path, {"file_type": "docx"})
    except Exception as e:
        logger.error(f"【WORD文件加载】加载文件 {abs_file_path} 时出错: {e}")
        return []


def markdown_loader_sync(file_path: str) -> list[Document]:
    """
    同步加载Markdown文件内容（用于多线程场景）
    :param file_path: Markdown文件路径
    :return: Markdown文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        return _single_document(_read_text_file(abs_file_path), abs_file_path, {"file_type": "markdown"})
    except Exception as e:
        logger.error(f"【Markdown文件加载】加载文件 {abs_file_path} 时出错: {e}")
        return []


def ppt_loader_sync(file_path: str) -> list[Document]:
    """
    同步加载PPT/PPTX文件内容（用于多线程场景）
    :param file_path: PPT文件路径
    :return: PPT文件内容
    """
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    lower_path = abs_file_path.lower()
    if lower_path.endswith(".ppt") and not lower_path.endswith(".pptx"):
        logger.error("【PPT文件加载】旧版 .ppt 文件需要先另存为 .pptx 后上传")
        return []

    try:
        return _extract_pptx_documents(abs_file_path)
    except Exception as e:
        logger.error(f"【PPT文件加载】加载文件 {abs_file_path} 时出错: {e}")
        return []


def _read_text_file(file_path: str) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            with open(file_path, encoding=encoding) as handle:
                return handle.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, encoding="utf-8", errors="replace") as handle:
        return handle.read()


def _single_document(content: str, source: str, metadata: dict | None = None) -> list[Document]:
    normalized = content.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not normalized:
        return []
    return [Document(page_content=normalized, metadata={"source": source, **(metadata or {})})]


def _iter_docx_blocks(document):
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _docx_table_text(table) -> str:
    rows = []
    for row in table.rows:
        cells = [" ".join(cell.text.split()) for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _extract_docx_text(file_path: str) -> str:
    from docx import Document as DocxDocument

    document = DocxDocument(file_path)
    parts: list[str] = []
    for block in _iter_docx_blocks(document):
        if hasattr(block, "rows"):
            text = _docx_table_text(block)
        else:
            text = block.text.strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _pptx_shape_lines(shape) -> list[str]:
    lines: list[str] = []

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
            if text:
                lines.append(text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [" ".join(cell.text.split()) for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))

    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            lines.extend(_pptx_shape_lines(child))

    return lines


def _extract_pptx_documents(file_path: str) -> list[Document]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    documents: list[Document] = []
    for index, slide in enumerate(presentation.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_pptx_shape_lines(shape))
        content = "\n".join(dict.fromkeys(line for line in lines if line))
        if content.strip():
            documents.append(
                Document(
                    page_content=f"幻灯片 {index}\n\n{content}",
                    metadata={"source": file_path, "page": index, "slide_number": index, "file_type": "pptx"},
                )
            )
    return documents
