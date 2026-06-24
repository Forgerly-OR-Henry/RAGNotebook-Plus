"""
模块职责：业务服务模块，负责组织领域用例、数据访问和外部能力协作。

主要协作：本文件只声明当前模块的职责边界，运行时行为由下方函数、类和依赖对象共同完成。
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import html
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from core.logger_handler import logger
from utils.env_loader import optional_env_value
from utils.path_tool import get_data_path


OFFICE_EXTENSIONS = {".doc", ".docx", ".ppt", ".pptx"}
OOXML_FALLBACK_EXTENSIONS = {".docx", ".pptx"}
EMU_PER_INCH = 914400
PX_PER_INCH = 96


class DocumentPreviewUnavailable(RuntimeError):
    """Raised when a document cannot be rendered in the current environment."""


@dataclass(frozen=True)
class RenderedDocumentPreview:
    """
    用途：领域对象或协作组件，用于承载本模块内的核心状态和行为。

    属性：
    - content（bytes）：保存content相关状态、配置或数据字段。
    - media_type（str）：保存media_type相关状态、配置或数据字段。
    - filename（str）：保存filename相关状态、配置或数据字段。
    - renderer（str）：保存renderer相关状态、配置或数据字段。
    """
    content: bytes
    media_type: str
    filename: str
    renderer: str


class DocumentPreviewService:
    """Render knowledge documents for visual preview.

    Office formats first try a real document renderer (LibreOffice/soffice) and
    fall back to self-contained HTML for modern OOXML files when that binary is
    not available.
    """

    def __init__(self, cache_dir: Path | None = None):
        """
        用途：执行init相关业务逻辑。

        参数：
        - cache_dir（Path | None）：调用方传入的cache_dir数据或控制参数，用于驱动本函数处理流程。

        返回：未显式标注；返回值供调用方继续编排业务流程或生成接口响应。
        """
        self.cache_dir = cache_dir or Path(get_data_path()) / "preview_cache"

    async def render(
        self,
        *,
        filename: str | None,
        file_ext: str | None,
        mime_type: str | None,
        content: bytes,
        content_hash: str | None,
    ) -> RenderedDocumentPreview:
        """
        用途：异步执行render相关业务流程。

        参数：
        - filename（str | None）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。
        - file_ext（str | None）：调用方传入的file_ext数据或控制参数，用于驱动本函数处理流程。
        - mime_type（str | None）：调用方传入的mime_type数据或控制参数，用于驱动本函数处理流程。
        - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
        - content_hash（str | None）：调用方传入的content_hash数据或控制参数，用于驱动本函数处理流程。

        返回：RenderedDocumentPreview；返回值供调用方继续编排业务流程或生成接口响应。

        副作用：可能访问数据库、文件、模型服务或流式事件通道，异常会沿调用链抛出。
        """
        return await asyncio.to_thread(
            self.render_sync,
            filename=filename,
            file_ext=file_ext,
            mime_type=mime_type,
            content=content,
            content_hash=content_hash,
        )

    def render_sync(
        self,
        *,
        filename: str | None,
        file_ext: str | None,
        mime_type: str | None,
        content: bytes,
        content_hash: str | None,
    ) -> RenderedDocumentPreview:
        """
        用途：渲染render sync相关的数据或流程。

        参数：
        - filename（str | None）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。
        - file_ext（str | None）：调用方传入的file_ext数据或控制参数，用于驱动本函数处理流程。
        - mime_type（str | None）：调用方传入的mime_type数据或控制参数，用于驱动本函数处理流程。
        - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
        - content_hash（str | None）：调用方传入的content_hash数据或控制参数，用于驱动本函数处理流程。

        返回：RenderedDocumentPreview；返回值供调用方继续编排业务流程或生成接口响应。
        """
        ext = _normalize_extension(file_ext, filename, mime_type)
        preview_name = _preview_filename(filename, ext)

        if ext == ".pdf":
            return RenderedDocumentPreview(
                content=content,
                media_type="application/pdf",
                filename=f"{Path(preview_name).stem}.pdf",
                renderer="native-pdf",
            )

        if ext not in OFFICE_EXTENSIONS:
            raise DocumentPreviewUnavailable("当前文件格式暂无原样预览器")

        converted = self._convert_office_to_pdf(
            content=content,
            content_hash=content_hash,
            source_ext=ext,
            preview_name=preview_name,
        )
        if converted is not None:
            return converted

        if ext == ".docx":
            return RenderedDocumentPreview(
                content=render_docx_html(content, filename or preview_name),
                media_type="text/html; charset=utf-8",
                filename=f"{Path(preview_name).stem}.html",
                renderer="docx-html-fallback",
            )
        if ext == ".pptx":
            return RenderedDocumentPreview(
                content=render_pptx_html(content, filename or preview_name),
                media_type="text/html; charset=utf-8",
                filename=f"{Path(preview_name).stem}.html",
                renderer="pptx-html-fallback",
            )

        raise DocumentPreviewUnavailable(
            "当前环境未安装 LibreOffice/soffice 转换器，旧版 .doc/.ppt 无法原样预览；请安装 LibreOffice 或将文件另存为 .docx/.pptx"
        )

    def _convert_office_to_pdf(
        self,
        *,
        content: bytes,
        content_hash: str | None,
        source_ext: str,
        preview_name: str,
    ) -> RenderedDocumentPreview | None:
        """
        用途：执行convert office to pdf相关业务逻辑。

        参数：
        - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
        - content_hash（str | None）：调用方传入的content_hash数据或控制参数，用于驱动本函数处理流程。
        - source_ext（str）：调用方传入的source_ext数据或控制参数，用于驱动本函数处理流程。
        - preview_name（str）：调用方传入的preview_name数据或控制参数，用于驱动本函数处理流程。

        返回：RenderedDocumentPreview | None；返回值供调用方继续编排业务流程或生成接口响应。
        """
        converter = _find_soffice()
        if not converter:
            return None

        digest = content_hash or hashlib.sha256(content).hexdigest()
        cache_key = f"{digest}-{source_ext.lstrip('.')}.pdf"
        cache_path = self.cache_dir / cache_key
        if cache_path.is_file():
            return RenderedDocumentPreview(
                content=cache_path.read_bytes(),
                media_type="application/pdf",
                filename=f"{Path(preview_name).stem}.pdf",
                renderer="libreoffice-cache",
            )

        try:
            self.cache_dir.mkdir(parents=True, exist_ok=True)
            with tempfile.TemporaryDirectory(prefix="ragnotebook-preview-") as temp_dir:
                temp_path = Path(temp_dir)
                input_path = temp_path / f"source{source_ext}"
                output_dir = temp_path / "out"
                profile_dir = temp_path / "lo-profile"
                output_dir.mkdir()
                profile_dir.mkdir()
                input_path.write_bytes(content)

                cmd = [
                    converter,
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    f"-env:UserInstallation={profile_dir.as_uri()}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(output_dir),
                    str(input_path),
                ]
                creationflags = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0
                completed = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=90,
                    creationflags=creationflags,
                    check=False,
                )
                if completed.returncode != 0:
                    logger.warning(
                        "LibreOffice preview conversion failed: %s %s",
                        completed.stdout.strip(),
                        completed.stderr.strip(),
                    )
                    return None

                pdf_files = sorted(output_dir.glob("*.pdf"))
                if not pdf_files:
                    logger.warning("LibreOffice preview conversion produced no PDF")
                    return None

                rendered = pdf_files[0].read_bytes()
                cache_path.write_bytes(rendered)
                return RenderedDocumentPreview(
                    content=rendered,
                    media_type="application/pdf",
                    filename=f"{Path(preview_name).stem}.pdf",
                    renderer="libreoffice",
                )
        except (OSError, subprocess.SubprocessError, TimeoutError) as exc:
            logger.warning("Office preview conversion unavailable: %s", exc)
            return None


def _normalize_extension(file_ext: str | None, filename: str | None, mime_type: str | None) -> str:
    """
    用途：执行normalize extension相关业务逻辑。

    参数：
    - file_ext（str | None）：调用方传入的file_ext数据或控制参数，用于驱动本函数处理流程。
    - filename（str | None）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。
    - mime_type（str | None）：调用方传入的mime_type数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    ext = (file_ext or Path(filename or "").suffix or "").strip().lower()
    if ext and not ext.startswith("."):
        ext = f".{ext}"
    if ext:
        return ext
    if mime_type == "application/pdf":
        return ".pdf"
    return ""


def _preview_filename(filename: str | None, ext: str) -> str:
    """
    用途：执行preview filename相关业务逻辑。

    参数：
    - filename（str | None）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。
    - ext（str）：调用方传入的ext数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    stem = Path(filename or "document").stem or "document"
    safe_stem = re.sub(r'[\\/:*?"<>|]+', "_", stem).strip() or "document"
    return f"{safe_stem}{ext or ''}"


def _find_soffice() -> str | None:
    """
    用途：执行find soffice相关业务逻辑。

    参数：无显式业务参数。

    返回：str | None；返回值供调用方继续编排业务流程或生成接口响应。
    """
    candidates = [
        optional_env_value("SOFFICE_PATH"),
        optional_env_value("LIBREOFFICE_PATH"),
        "soffice",
        "libreoffice",
    ]
    if os.name == "nt":
        candidates.extend(
            [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\Program Files\OpenOffice 4\program\soffice.exe",
                r"C:\Program Files (x86)\OpenOffice 4\program\soffice.exe",
            ]
        )

    for candidate in candidates:
        if not candidate:
            continue
        path = Path(candidate)
        if path.is_file():
            return str(path)
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def render_docx_html(content: bytes, filename: str) -> bytes:
    """
    用途：渲染render docx html相关的数据或流程。

    参数：
    - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
    - filename（str）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。

    返回：bytes；返回值供调用方继续编排业务流程或生成接口响应。
    """
    from docx import Document as DocxDocument
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = DocxDocument(BytesIO(content))
    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(_render_docx_paragraph(Paragraph(child, document)))
        elif isinstance(child, CT_Tbl):
            blocks.append(_render_docx_table(Table(child, document)))

    if not any(block.strip() for block in blocks):
        blocks.append('<p class="docx-empty">暂无可预览内容</p>')

    title = html.escape(filename)
    body = "\n".join(blocks)
    return _wrap_preview_html(
        title=title,
        css=_DOCX_PREVIEW_CSS,
        body=f'<main class="docx-page" aria-label="{title}">{body}</main>',
    )


def _render_docx_paragraph(paragraph, compact: bool = False) -> str:
    """
    用途：渲染render docx paragraph相关的数据或流程。

    参数：
    - paragraph（未显式标注）：调用方传入的paragraph数据或控制参数，用于驱动本函数处理流程。
    - compact（bool）：调用方传入的compact数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    content = "".join(_render_docx_run(run, paragraph) for run in paragraph.runs)
    if not content and paragraph.text:
        content = _escape_inline_text(paragraph.text)
    if not content:
        content = "&nbsp;"

    style_name = (getattr(paragraph.style, "name", "") or "").lower()
    tag = "p"
    class_names = ["docx-paragraph"]
    if compact:
        class_names.append("docx-paragraph--compact")
    heading_match = re.match(r"heading\s+([1-6])", style_name)
    if heading_match:
        tag = f"h{heading_match.group(1)}"
        class_names.append("docx-heading")
    if _docx_is_numbered_or_bulleted(paragraph):
        class_names.append("docx-list")

    styles: list[str] = []
    alignment = _alignment_css(paragraph.alignment)
    if alignment:
        styles.append(f"text-align: {alignment}")
    fmt = paragraph.paragraph_format
    if fmt.left_indent:
        styles.append(f"margin-left: {_length_to_pt(fmt.left_indent)}pt")
    if fmt.first_line_indent:
        styles.append(f"text-indent: {_length_to_pt(fmt.first_line_indent)}pt")
    if fmt.space_before:
        styles.append(f"margin-top: {_length_to_pt(fmt.space_before)}pt")
    if fmt.space_after:
        styles.append(f"margin-bottom: {_length_to_pt(fmt.space_after)}pt")
    if fmt.line_spacing and isinstance(fmt.line_spacing, float):
        styles.append(f"line-height: {fmt.line_spacing:.3g}")

    attrs = _html_attrs({"class": " ".join(class_names), "style": "; ".join(styles)})
    return f"<{tag}{attrs}>{content}</{tag}>"


def _render_docx_run(run, paragraph) -> str:
    """
    用途：渲染render docx run相关的数据或流程。

    参数：
    - run（未显式标注）：调用方传入的run数据或控制参数，用于驱动本函数处理流程。
    - paragraph（未显式标注）：调用方传入的paragraph数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    fragments: list[str] = []
    if run.text:
        fragments.append(_escape_inline_text(run.text))
    fragments.extend(_docx_run_images(run))
    if not fragments:
        return ""

    styles: list[str] = []
    size = run.font.size or getattr(paragraph.style.font, "size", None)
    if size:
        styles.append(f"font-size: {size.pt:.2f}pt")
    color = _docx_color(run.font.color) or _docx_color(getattr(paragraph.style.font, "color", None))
    if color:
        styles.append(f"color: {color}")
    if run.bold:
        styles.append("font-weight: 700")
    if run.italic:
        styles.append("font-style: italic")
    if run.underline:
        styles.append("text-decoration: underline")

    attrs = _html_attrs({"style": "; ".join(styles)})
    return f"<span{attrs}>{''.join(fragments)}</span>"


def _docx_run_images(run) -> list[str]:
    """
    用途：执行docx run images相关业务逻辑。

    参数：
    - run（未显式标注）：调用方传入的run数据或控制参数，用于驱动本函数处理流程。

    返回：list[str]；返回值供调用方继续编排业务流程或生成接口响应。
    """
    from docx.oxml.ns import qn

    images: list[str] = []
    blips = run._element.xpath(".//*[local-name()='blip']")
    extents = run._element.xpath(".//*[local-name()='extent']")
    for index, blip in enumerate(blips):
        rel_id = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
        if not rel_id:
            continue
        image_part = run.part.related_parts.get(rel_id)
        if not image_part or not getattr(image_part, "blob", None):
            continue
        data_url = _data_url(image_part.blob, getattr(image_part, "content_type", None))
        styles = ["max-width: 100%", "height: auto", "vertical-align: middle"]
        if index < len(extents):
            cx = _int_attr(extents[index].get("cx"))
            cy = _int_attr(extents[index].get("cy"))
            if cx:
                styles.append(f"width: {_emu_to_px(cx):.2f}px")
            if cy:
                styles.append(f"max-height: {_emu_to_px(cy):.2f}px")
        images.append(f'<img class="docx-image" src="{data_url}" alt="" style="{"; ".join(styles)}" />')
    return images


def _render_docx_table(table) -> str:
    """
    用途：渲染render docx table相关的数据或流程。

    参数：
    - table（未显式标注）：调用方传入的table数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    rows: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_body = "".join(_render_docx_paragraph(paragraph, compact=True) for paragraph in cell.paragraphs)
            for nested_table in cell.tables:
                cell_body += _render_docx_table(nested_table)
            attrs = _html_attrs({"style": _docx_cell_style(cell)})
            cells.append(f"<td{attrs}>{cell_body or '&nbsp;'}</td>")
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f'<table class="docx-table"><tbody>{"".join(rows)}</tbody></table>'


def _docx_cell_style(cell) -> str:
    """
    用途：执行docx cell style相关业务逻辑。

    参数：
    - cell（未显式标注）：调用方传入的cell数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    from docx.oxml.ns import qn

    styles: list[str] = []
    shading = cell._tc.xpath(".//*[local-name()='shd']")
    if shading:
        fill = shading[0].get(qn("w:fill"))
        if fill and fill.lower() != "auto":
            styles.append(f"background: #{fill}")
    return "; ".join(styles)


def _docx_is_numbered_or_bulleted(paragraph) -> bool:
    """
    用途：执行docx is numbered or bulleted相关业务逻辑。

    参数：
    - paragraph（未显式标注）：调用方传入的paragraph数据或控制参数，用于驱动本函数处理流程。

    返回：bool；返回值供调用方继续编排业务流程或生成接口响应。
    """
    p_pr = paragraph._p.pPr
    return bool(p_pr is not None and p_pr.numPr is not None)


def render_pptx_html(content: bytes, filename: str) -> bytes:
    """
    用途：渲染render pptx html相关的数据或流程。

    参数：
    - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
    - filename（str）：调用方传入的filename数据或控制参数，用于驱动本函数处理流程。

    返回：bytes；返回值供调用方继续编排业务流程或生成接口响应。
    """
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        shapes = "".join(_render_pptx_shape(shape, slide_width, slide_height) for shape in slide.shapes)
        background = _pptx_slide_background(slide)
        style = f"--slide-w: {slide_width}; --slide-h: {slide_height}; background: {background}"
        slides.append(
            f'<section class="ppt-slide" aria-label="Slide {index}" style="{style}">'
            f'<div class="ppt-slide-number">{index}</div>{shapes}</section>'
        )

    if not slides:
        slides.append('<section class="ppt-slide ppt-slide--empty">暂无可预览内容</section>')

    title = html.escape(filename)
    return _wrap_preview_html(
        title=title,
        css=_PPTX_PREVIEW_CSS,
        body=f'<main class="ppt-stage" aria-label="{title}">{"".join(slides)}</main>',
    )


def _render_pptx_shape(shape, slide_width: int, slide_height: int) -> str:
    """
    用途：渲染render pptx shape相关的数据或流程。

    参数：
    - shape（未显式标注）：调用方传入的shape数据或控制参数，用于驱动本函数处理流程。
    - slide_width（int）：调用方传入的slide_width数据或控制参数，用于驱动本函数处理流程。
    - slide_height（int）：调用方传入的slide_height数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    if hasattr(shape, "shapes"):
        return "".join(_render_pptx_shape(child, slide_width, slide_height) for child in shape.shapes)

    box_style = _pptx_box_style(shape, slide_width, slide_height)
    fill_style = _pptx_fill_style(shape)
    line_style = _pptx_line_style(shape)
    styles = "; ".join(part for part in [box_style, fill_style, line_style] if part)

    if _pptx_shape_has_image(shape):
        image = shape.image
        data_url = _data_url(image.blob, image.content_type)
        return f'<img class="ppt-shape ppt-image" src="{data_url}" alt="" style="{styles}" />'

    if getattr(shape, "has_table", False):
        table = _render_pptx_table(shape.table)
        return f'<div class="ppt-shape ppt-table-wrap" style="{styles}">{table}</div>'

    if getattr(shape, "has_text_frame", False):
        text = _render_pptx_text_frame(shape.text_frame)
        if text.strip():
            return f'<div class="ppt-shape ppt-textbox" style="{styles}">{text}</div>'

    if fill_style or line_style:
        return f'<div class="ppt-shape ppt-geometry" style="{styles}"></div>'
    return ""


def _pptx_box_style(shape, slide_width: int, slide_height: int) -> str:
    """
    用途：执行pptx box style相关业务逻辑。

    参数：
    - shape（未显式标注）：调用方传入的shape数据或控制参数，用于驱动本函数处理流程。
    - slide_width（int）：调用方传入的slide_width数据或控制参数，用于驱动本函数处理流程。
    - slide_height（int）：调用方传入的slide_height数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    styles = [
        "position: absolute",
        f"left: {_ratio_pct(shape.left, slide_width)}",
        f"top: {_ratio_pct(shape.top, slide_height)}",
        f"width: {_ratio_pct(shape.width, slide_width)}",
        f"height: {_ratio_pct(shape.height, slide_height)}",
    ]
    rotation = float(getattr(shape, "rotation", 0) or 0)
    if rotation:
        styles.append(f"transform: rotate({rotation:.2f}deg)")
        styles.append("transform-origin: center")
    return "; ".join(styles)


def _render_pptx_text_frame(text_frame) -> str:
    """
    用途：渲染render pptx text frame相关的数据或流程。

    参数：
    - text_frame（未显式标注）：调用方传入的text_frame数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        runs = "".join(_render_pptx_run(run) for run in paragraph.runs)
        if not runs and paragraph.text:
            runs = _escape_inline_text(paragraph.text)
        if not runs:
            continue
        styles: list[str] = []
        alignment = _alignment_css(paragraph.alignment)
        if alignment:
            styles.append(f"text-align: {alignment}")
        margin_left = max(int(getattr(paragraph, "level", 0) or 0), 0) * 1.1
        if margin_left:
            styles.append(f"margin-left: {margin_left:.1f}em")
        attrs = _html_attrs({"style": "; ".join(styles)})
        bullet = '<span class="ppt-bullet-dot">•</span>' if margin_left else ""
        paragraphs.append(f"<p{attrs}>{bullet}{runs}</p>")
    return f'<div class="ppt-text">{"".join(paragraphs)}</div>'


def _render_pptx_run(run) -> str:
    """
    用途：渲染render pptx run相关的数据或流程。

    参数：
    - run（未显式标注）：调用方传入的run数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    text = _escape_inline_text(run.text)
    if not text:
        return ""
    styles: list[str] = []
    size = getattr(run.font, "size", None)
    if size:
        styles.append(f"font-size: {size.pt:.2f}pt")
    color = _pptx_color(getattr(run.font, "color", None))
    if color:
        styles.append(f"color: {color}")
    if run.font.bold:
        styles.append("font-weight: 700")
    if run.font.italic:
        styles.append("font-style: italic")
    if run.font.underline:
        styles.append("text-decoration: underline")
    attrs = _html_attrs({"style": "; ".join(styles)})
    return f"<span{attrs}>{text}</span>"


def _render_pptx_table(table) -> str:
    """
    用途：渲染render pptx table相关的数据或流程。

    参数：
    - table（未显式标注）：调用方传入的table数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    rows: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            styles: list[str] = []
            try:
                fill = _pptx_color(cell.fill.fore_color)
            except Exception:
                fill = None
            if fill:
                styles.append(f"background: {fill}")
            attrs = _html_attrs({"style": "; ".join(styles)})
            cells.append(f"<td{attrs}>{_render_pptx_text_frame(cell.text_frame) or '&nbsp;'}</td>")
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f'<table class="ppt-table"><tbody>{"".join(rows)}</tbody></table>'


def _pptx_shape_has_image(shape) -> bool:
    """
    用途：执行pptx shape has image相关业务逻辑。

    参数：
    - shape（未显式标注）：调用方传入的shape数据或控制参数，用于驱动本函数处理流程。

    返回：bool；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        _ = shape.image
        return True
    except Exception:
        return False


def _pptx_slide_background(slide) -> str:
    """
    用途：执行pptx slide background相关业务逻辑。

    参数：
    - slide（未显式标注）：调用方传入的slide数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        color = _pptx_color(slide.background.fill.fore_color)
        if color:
            return color
    except Exception:
        pass
    return "#ffffff"


def _pptx_fill_style(shape) -> str:
    """
    用途：执行pptx fill style相关业务逻辑。

    参数：
    - shape（未显式标注）：调用方传入的shape数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        color = _pptx_color(shape.fill.fore_color)
        if color:
            return f"background: {color}"
    except Exception:
        pass
    return ""


def _pptx_line_style(shape) -> str:
    """
    用途：执行pptx line style相关业务逻辑。

    参数：
    - shape（未显式标注）：调用方传入的shape数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        color = _pptx_color(shape.line.color)
        width = getattr(shape.line, "width", None)
        if color:
            border_width = f"{max(_emu_to_px(int(width)), 1):.2f}px" if width else "1px"
            return f"border: {border_width} solid {color}"
    except Exception:
        pass
    return ""


def _pptx_color(color_format) -> str | None:
    """
    用途：执行pptx color相关业务逻辑。

    参数：
    - color_format（未显式标注）：调用方传入的color_format数据或控制参数，用于驱动本函数处理流程。

    返回：str | None；返回值供调用方继续编排业务流程或生成接口响应。
    """
    if color_format is None:
        return None
    try:
        rgb = color_format.rgb
    except Exception:
        return None
    return f"#{rgb}" if rgb else None


def _docx_color(color_format) -> str | None:
    """
    用途：执行docx color相关业务逻辑。

    参数：
    - color_format（未显式标注）：调用方传入的color_format数据或控制参数，用于驱动本函数处理流程。

    返回：str | None；返回值供调用方继续编排业务流程或生成接口响应。
    """
    if color_format is None:
        return None
    try:
        rgb = color_format.rgb
    except Exception:
        return None
    return f"#{rgb}" if rgb else None


def _alignment_css(value) -> str | None:
    """
    用途：执行alignment css相关业务逻辑。

    参数：
    - value（未显式标注）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。

    返回：str | None；返回值供调用方继续编排业务流程或生成接口响应。
    """
    name = (getattr(value, "name", "") or "").lower()
    if "center" in name:
        return "center"
    if "right" in name:
        return "right"
    if "justify" in name or "distribute" in name:
        return "justify"
    if "left" in name:
        return "left"
    return None


def _escape_inline_text(value: str) -> str:
    """
    用途：执行escape inline text相关业务逻辑。

    参数：
    - value（str）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    return html.escape(value).replace("\t", "&emsp;").replace("\n", "<br>")


def _html_attrs(values: dict[str, str | None]) -> str:
    """
    用途：执行html attrs相关业务逻辑。

    参数：
    - values（dict[str, str | None]）：调用方传入的values数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    attrs = []
    for key, value in values.items():
        if value:
            attrs.append(f' {key}="{html.escape(value, quote=True)}"')
    return "".join(attrs)


def _data_url(content: bytes, content_type: str | None) -> str:
    """
    用途：执行data url相关业务逻辑。

    参数：
    - content（bytes）：调用方传入的content数据或控制参数，用于驱动本函数处理流程。
    - content_type（str | None）：调用方传入的content_type数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    media_type = content_type or "application/octet-stream"
    encoded = base64.b64encode(content).decode("ascii")
    return f"data:{media_type};base64,{encoded}"


def _length_to_pt(value) -> float:
    """
    用途：执行length to pt相关业务逻辑。

    参数：
    - value（未显式标注）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。

    返回：float；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        return float(value.pt)
    except Exception:
        return 0.0


def _emu_to_px(value: int) -> float:
    """
    用途：执行emu to px相关业务逻辑。

    参数：
    - value（int）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。

    返回：float；返回值供调用方继续编排业务流程或生成接口响应。
    """
    return value / EMU_PER_INCH * PX_PER_INCH


def _ratio_pct(value: int, total: int) -> str:
    """
    用途：执行ratio pct相关业务逻辑。

    参数：
    - value（int）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。
    - total（int）：调用方传入的total数据或控制参数，用于驱动本函数处理流程。

    返回：str；返回值供调用方继续编排业务流程或生成接口响应。
    """
    if not total:
        return "0%"
    return f"{value / total * 100:.5f}%"


def _int_attr(value: str | None) -> int | None:
    """
    用途：执行int attr相关业务逻辑。

    参数：
    - value（str | None）：调用方传入的value数据或控制参数，用于驱动本函数处理流程。

    返回：int | None；返回值供调用方继续编排业务流程或生成接口响应。
    """
    try:
        return int(value) if value is not None else None
    except ValueError:
        return None


def _wrap_preview_html(*, title: str, css: str, body: str) -> bytes:
    """
    用途：执行wrap preview html相关业务逻辑。

    参数：
    - title（str）：调用方传入的title数据或控制参数，用于驱动本函数处理流程。
    - css（str）：调用方传入的css数据或控制参数，用于驱动本函数处理流程。
    - body（str）：调用方传入的body数据或控制参数，用于驱动本函数处理流程。

    返回：bytes；返回值供调用方继续编排业务流程或生成接口响应。
    """
    document = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{title}</title>
  <style>{css}</style>
</head>
<body>
{body}
</body>
</html>
"""
    return document.encode("utf-8")


_DOCX_PREVIEW_CSS = """
html, body {
  margin: 0;
  background: #e8ecf2;
  color: #111827;
  font-family: Arial, "Microsoft YaHei", sans-serif;
}
body {
  padding: 24px;
}
.docx-page {
  box-sizing: border-box;
  width: min(900px, calc(100vw - 48px));
  min-height: 1120px;
  margin: 0 auto;
  padding: 72px 80px;
  background: #fff;
  box-shadow: 0 14px 40px rgba(15, 23, 42, 0.16);
}
.docx-paragraph {
  margin: 0 0 10pt;
  line-height: 1.5;
  white-space: pre-wrap;
}
.docx-paragraph--compact {
  margin-bottom: 4pt;
}
.docx-heading {
  margin: 0 0 12pt;
  line-height: 1.25;
  font-weight: 700;
}
h1.docx-heading { font-size: 24pt; }
h2.docx-heading { font-size: 20pt; }
h3.docx-heading { font-size: 17pt; }
h4.docx-heading { font-size: 15pt; }
h5.docx-heading, h6.docx-heading { font-size: 13pt; }
.docx-list::before {
  content: "•";
  display: inline-block;
  width: 1em;
  margin-left: -1em;
}
.docx-table {
  width: 100%;
  margin: 12pt 0;
  border-collapse: collapse;
  table-layout: fixed;
}
.docx-table td,
.docx-table th {
  border: 1px solid #9ca3af;
  padding: 6pt 8pt;
  vertical-align: top;
}
.docx-image {
  display: inline-block;
  margin: 4pt 0;
}
@media (max-width: 720px) {
  body { padding: 12px; }
  .docx-page {
    width: 100%;
    min-height: 0;
    padding: 34px 24px;
  }
}
"""


_PPTX_PREVIEW_CSS = """
html, body {
  margin: 0;
  background: #111827;
  color: #111827;
  font-family: Arial, "Microsoft YaHei", sans-serif;
}
.ppt-stage {
  display: grid;
  justify-items: center;
  gap: 24px;
  padding: 24px;
}
.ppt-slide {
  position: relative;
  width: min(960px, calc(100vw - 48px));
  aspect-ratio: var(--slide-w) / var(--slide-h);
  overflow: hidden;
  box-shadow: 0 18px 48px rgba(0, 0, 0, 0.34);
}
.ppt-shape {
  box-sizing: border-box;
  overflow: hidden;
}
.ppt-image {
  object-fit: contain;
}
.ppt-textbox {
  padding: 0.35em 0.45em;
}
.ppt-text {
  width: 100%;
  height: 100%;
}
.ppt-text p {
  margin: 0;
  line-height: 1.2;
  white-space: pre-wrap;
}
.ppt-bullet-dot {
  display: inline-block;
  width: 1em;
  margin-left: -1em;
}
.ppt-table {
  width: 100%;
  height: 100%;
  border-collapse: collapse;
  table-layout: fixed;
  font-size: 12px;
}
.ppt-table td,
.ppt-table th {
  border: 1px solid #6b7280;
  padding: 4px 6px;
  vertical-align: top;
}
.ppt-geometry {
  min-width: 1px;
  min-height: 1px;
}
.ppt-slide-number {
  position: absolute;
  right: 10px;
  bottom: 8px;
  z-index: 20;
  color: rgba(17, 24, 39, 0.45);
  font-size: 11px;
}
@media (max-width: 720px) {
  .ppt-stage {
    gap: 14px;
    padding: 12px;
  }
  .ppt-slide {
    width: calc(100vw - 24px);
  }
}
"""
